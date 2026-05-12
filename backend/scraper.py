import requests
import trafilatura
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse
import docx
import PyPDF2
from pptx import Presentation
import socket
import ipaddress
import re

def is_safe_url(url):
    try:
        hostname = urlparse(url).hostname
        if not hostname: return False
        
        domain = hostname.lower()
        if domain != 'smkn2indramayu.sch.id' and not domain.endswith('.smkn2indramayu.sch.id'):
            return False
            
        ip = socket.gethostbyname(hostname)
        ip_obj = ipaddress.ip_address(ip)
        if ip_obj.is_private or ip_obj.is_loopback or ip_obj.is_link_local:
            return False
        return True
    except Exception:
        return False

HEADERS = {
    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
}

def extract_text_from_pdf(file_stream):
    """Mengekstrak teks dari file PDF."""
    try:
        reader = PyPDF2.PdfReader(file_stream)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return None

def extract_text_from_docx(file_stream):
    """Mengekstrak teks dari file DOCX."""
    try:
        doc = docx.Document(file_stream)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"Error reading DOCX: {e}")
        return None

def extract_text_from_pptx(file_stream):
    """Mengekstrak teks dari file PPTX."""
    try:
        prs = Presentation(file_stream)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return None


def clean_html_boilerplate(html_text):
    soup = BeautifulSoup(html_text, 'html.parser')
    for tag in soup.find_all(['nav', 'footer', 'header', 'aside', 'script', 'style', 'noscript']):
        tag.decompose()
    junk_selectors = [
        '.sidebar', '#sidebar', '.menu', '#menu', '.navbar', '#navbar', 
        '.widget', '.footer', '.comments', '.ad', '.advertisement', 
        '.cookie-banner', '.popup'
    ]
    for selector in junk_selectors:
        for element in soup.select(selector):
            element.decompose()
    return str(soup)

def extract_single_page(url):
    """
    Extracts content and the single best thumbnail image URL from a page.
    FIX #5: Stores only ONE image_url (og:image preferred, else first in-content image).
    This prevents comma-joined multi-URL strings from breaking img src attributes.
    """
    try:
        if not is_safe_url(url):
            return {"status": "error", "url": url, "reason": "URL is not allowed (SSRF protection).", "image_url": None}
            
        response = requests.get(url, headers=HEADERS, timeout=15)
        response.raise_for_status()
        response.encoding = response.apparent_encoding
        
        html_content = response.text
        cleaned_html = clean_html_boilerplate(html_content)
        content = trafilatura.extract(
            cleaned_html, 
            include_comments=False, 
            include_tables=True,
            favor_precision=True,
            target_language="id"
        )
        title = trafilatura.extract_metadata(html_content).title if trafilatura.extract_metadata(html_content) else ""
        
        soup = BeautifulSoup(html_content, 'html.parser')
        primary_image_url = None

        # Priority 1: og:image — best thumbnail/representative image
        og_image = soup.find('meta', property='og:image')
        if og_image and og_image.get('content'):
            primary_image_url = urljoin(url, og_image['content'])

        # Priority 2: First <img> inside the extracted content body
        if not primary_image_url and content:
            content_soup = BeautifulSoup(content, 'html.parser')
            first_img = content_soup.find('img')
            if first_img and first_img.get('src'):
                primary_image_url = urljoin(url, first_img['src'])

        if content:
            content = re.sub(r'\n\s*\n', '\n\n', content).strip()
            if len(content) < 150:
                return {"status": "skipped", "url": url, "reason": "Content too short or mostly boilerplate", "image_url": None}
            return {"status": "success", "url": url, "title": title, "content": content, "image_url": primary_image_url}
        else:
            return {"status": "skipped", "url": url, "reason": "No main content found", "image_url": None}

    except requests.exceptions.RequestException as e:
        return {"status": "error", "url": url, "reason": str(e), "image_url": None}

def scrape_from_file(file_path):
    """
    Reads a file of URLs and yields the extraction result for each URL.
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            urls = [line.strip() for line in f if line.strip() and not line.startswith('#')]
    except FileNotFoundError:
        yield {"status": "error", "url": file_path, "reason": "File not found"}
        return
        
    yield {"status": "info", "message": f"Ditemukan {len(urls)} URL untuk di-scrape."}

    for url in urls:
        yield extract_single_page(url)

def crawl_website(base_url, max_pages=50):
    """
    Crawls a website starting from base_url, discovering internal links, 
    and scraping each page up to max_pages.
    """
    yield {"status": "info", "message": f"Memulai deep crawl di: {base_url} (max {max_pages} halaman)"}
    
    domain = urlparse(base_url).netloc
    if not domain:
        yield {"status": "error", "url": base_url, "reason": "URL tidak valid."}
        return

    visited = set()
    queue = [base_url]
    
    # Simple heuristic to avoid downloading large binaries during link discovery
    ignored_extensions = {'.pdf', '.zip', '.rar', '.doc', '.docx', '.ppt', '.pptx', '.xls', '.xlsx', '.png', '.jpg', '.jpeg', '.gif', '.mp4', '.mp3'}
    ignored_paths = {'/login', '/admin', '/wp-admin', '/logout', '/tag/', '/category/', '/author/', '/page/'}
    
    while queue and len(visited) < max_pages:
        current_url = queue.pop(0)
        
        if current_url in visited:
            continue
            
        visited.add(current_url)
        yield {"status": "info", "message": f"[{len(visited)}/{max_pages}] Scrape: {current_url}"}
        
        # 1. Fetch the page to get links and content
        try:
            if not is_safe_url(current_url):
                yield {"status": "skipped", "url": current_url, "reason": "URL is not allowed (SSRF protection)."}
                continue
                
            response = requests.get(current_url, headers=HEADERS, timeout=15)
            # Only process if HTML
            if 'text/html' not in response.headers.get('Content-Type', ''):
                yield {"status": "skipped", "url": current_url, "reason": "Bukan HTML content"}
                continue
                
            response.raise_for_status()
            response.encoding = response.apparent_encoding
            html_content = response.text
            
            # 2. Find links for the queue
            soup = BeautifulSoup(html_content, 'html.parser')
            for a_tag in soup.find_all('a', href=True):
                href = a_tag['href']
                absolute_url = urljoin(current_url, href)
                # Remove fragments
                absolute_url = urlparse(absolute_url)._replace(fragment="").geturl()
                
                parsed_href = urlparse(absolute_url)
                
                if parsed_href.netloc == domain:
                    # Ignore specific extensions and paths
                    ext = parsed_href.path.lower()
                    if any(ext.endswith(e) for e in ignored_extensions): continue
                    if any(p in parsed_href.path.lower() for p in ignored_paths): continue
                        
                    if absolute_url not in visited and absolute_url not in queue:
                        queue.append(absolute_url)
                        
            # 3. Extract content (reusing logic from extract_single_page)
            cleaned_html = clean_html_boilerplate(html_content)
            content = trafilatura.extract(
                cleaned_html, 
                include_comments=False, 
                include_tables=True,
                favor_precision=True,
                target_language="id"
            )
            title = trafilatura.extract_metadata(html_content).title if trafilatura.extract_metadata(html_content) else ""
            
            primary_image_url = None
            og_image = soup.find('meta', property='og:image')
            if og_image and og_image.get('content'):
                primary_image_url = urljoin(current_url, og_image['content'])
            if not primary_image_url and content:
                content_soup = BeautifulSoup(content, 'html.parser')
                first_img = content_soup.find('img')
                if first_img and first_img.get('src'):
                    primary_image_url = urljoin(current_url, first_img['src'])

            if content:
                content = re.sub(r'\n\s*\n', '\n\n', content).strip()
                if len(content) < 150:
                    yield {"status": "skipped", "url": current_url, "reason": "Content too short or mostly boilerplate", "image_url": None}
                else:
                    yield {"status": "success", "url": current_url, "title": title, "content": content, "image_url": primary_image_url}
            else:
                yield {"status": "skipped", "url": current_url, "reason": "No main content found", "image_url": None}

        except Exception as e:
            yield {"status": "error", "url": current_url, "reason": str(e), "image_url": None}

    yield {"status": "info", "message": f"Crawl selesai. Total dikunjungi: {len(visited)} halaman."}